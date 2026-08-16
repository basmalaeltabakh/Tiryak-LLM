import fitz  # PyMuPDF
from pathlib import Path
from typing import List, Dict


def parse_pdf(file_path: str) -> List[Dict]:
    """
    Reads a PDF file and returns the content of each page.

    Returns:
        List of dictionaries in the format:
        {"page_number": int, "text": str}
    """
    doc = fitz.open(file_path)
    pages_content = []

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text")
        pages_content.append({
            "page_number": page_num,
            "text": text.strip()
        })

    doc.close()
    return pages_content


def is_scanned_pdf(pages_content: List[Dict], min_chars_threshold: int = 20) -> bool:
    """
    Determines whether a PDF is scanned or contains extractable text.

    If the average number of characters per page is below the threshold,
    the document is assumed to be scanned and requires OCR.
    """
    total_chars = sum(len(page["text"]) for page in pages_content)
    avg_chars_per_page = total_chars / len(pages_content) if pages_content else 0
    return avg_chars_per_page < min_chars_threshold


def parse_docx(file_path: str) -> List[Dict]:
    """
    Reads a DOCX file and returns its content as a single text block.

    Since DOCX files do not have fixed pages like PDFs,
    the entire document is returned as page 1.
    """
    from docx import Document

    doc = Document(file_path)
    full_text = []

    # Extract text from paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text.strip())

    # Extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            full_text.append(row_text)

    return [{"page_number": 1, "text": "\n".join(full_text)}]


def parse_pptx(file_path: str) -> List[Dict]:
    """
    Reads a PowerPoint file and returns the content of each slide.

    Each slide is treated as a separate page.
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    slides_content = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text.strip():
                        texts.append(text.strip())

        slides_content.append({
            "page_number": slide_num,
            "text": "\n".join(texts)
        })

    return slides_content


def parse_document(file_path: str) -> List[Dict]:
    """
    Unified entry point for document parsing.

    Detects the file type and uses the appropriate parser.
    Automatically falls back to OCR for scanned PDF files.
    """
    extension = Path(file_path).suffix.lower()

    if extension == ".pdf":
        pages_content = parse_pdf(file_path)

        # Use OCR if the PDF appears to be scanned
        if is_scanned_pdf(pages_content):
            from app.ingestion.ocr import ocr_scanned_pdf
            pages_content = ocr_scanned_pdf(file_path)

        return pages_content

    elif extension == ".docx":
        return parse_docx(file_path)

    elif extension == ".pptx":
        return parse_pptx(file_path)

    else:
        raise ValueError(f"Unsupported file format: {extension}")